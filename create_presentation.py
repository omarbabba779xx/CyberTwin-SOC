"""
CyberTwin SOC - Academic Defense (Soutenance) Presentation Generator
Generates a professional PowerPoint presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIGHTLY_LIGHTER = RGBColor(0x22, 0x22, 0x3e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x10, 0xB9, 0x81)
YELLOW = RGBColor(0xF5, 0xA6, 0x23)
DARK_ROW = RGBColor(0x16, 0x16, 0x2A)
HEADER_ROW = RGBColor(0x0A, 0x3D, 0x5C)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_title(slide, text, top=Inches(0.4), left=Inches(0.8), width=Inches(11.7), font_size=36):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    # underline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.85), Inches(2), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    return txBox


def add_bullet_slide(slide, title, bullets, top_start=Inches(1.7)):
    set_slide_bg(slide, BG_DARK)
    add_title(slide, title)
    txBox = slide.shapes.add_textbox(Inches(1.0), top_start, Inches(11.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)
        p.space_after = Pt(6)
        p.level = 0


def add_table(slide, rows, col_widths, left, top, row_height=Inches(0.5)):
    table_shape = slide.shapes.add_table(len(rows), len(col_widths), left, top,
                                          sum(col_widths), row_height * len(rows))
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if ri == 0:
                p.font.color.rgb = WHITE
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_ROW
            else:
                p.font.color.rgb = LIGHT_GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_ROW if ri % 2 == 1 else BG_SLIGHTLY_LIGHTER

            # remove cell borders
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is not None:
                    tcPr.remove(ln)

    return table


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════
    # SLIDE 1 - Title
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    # decorative top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CyberTwin SOC"
    p.font.size = Pt(54)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Jumeau Numerique pour la Simulation d'Attaques Cyber\net l'Evaluation de la Posture SOC"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # Author
    txBox3 = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.3), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "[Nom de l'Auteur]"
    p3.font.size = Pt(20)
    p3.font.color.rgb = PURPLE
    p3.alignment = PP_ALIGN.CENTER

    # Date
    txBox4 = slide.shapes.add_textbox(Inches(3), Inches(6.1), Inches(7.3), Inches(0.5))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "2026"
    p4.font.size = Pt(18)
    p4.font.color.rgb = LIGHT_GRAY
    p4.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 2 - Problematique
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    bullets = [
        "\u25b6  Les SOC font face a des menaces croissantes (APT, ransomware, insider threats)",
        "\u25b6  Comment tester la capacite de detection AVANT une vraie attaque ?",
        "\u25b6  Comment evaluer la maturite d'un SOC de maniere objective ?",
        "",
        "\u2192  Solution : Le Jumeau Numerique (Digital Twin) de cybersecurite",
    ]
    add_bullet_slide(slide, "Problematique", bullets)
    # Make solution line cyan
    for shp in slide.shapes:
        if shp.has_text_frame and len(shp.text_frame.paragraphs) >= 5:
            last_p = shp.text_frame.paragraphs[4]
            if last_p.text.strip().startswith("\u2192"):
                last_p.font.color.rgb = CYAN
                last_p.font.bold = True
                last_p.font.size = Pt(22)
                break

    # ════════════════════════════════════════════
    # SLIDE 3 - Objectifs
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Objectifs du Projet")

    objectives = [
        ("01", "Simuler un environnement IT realiste\n(jumeau numerique)", CYAN),
        ("02", "Reproduire des attaques APT documentees\n(APT29, APT28, TeamTNT)", PURPLE),
        ("03", "Evaluer la detection avec un scoring\nmulti-dimensionnel", GREEN),
        ("04", "Generer des rapports d'analyse AI\nautomatises", YELLOW),
    ]
    box_w = Inches(2.6)
    box_h = Inches(2.2)
    gap = Inches(0.3)
    start_x = Inches(0.7)
    y = Inches(2.0)

    for i, (num, text, accent) in enumerate(objectives):
        x = start_x + i * (box_w + gap)
        rect = add_shape_rect(slide, x, y, box_w, box_h, BG_SLIGHTLY_LIGHTER, accent)

        # Number
        numBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.6))
        ntf = numBox.text_frame
        np = ntf.paragraphs[0]
        np.text = num
        np.font.size = Pt(32)
        np.font.color.rgb = accent
        np.font.bold = True

        # Text
        txtBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), box_w - Inches(0.4), Inches(1.2))
        ttf = txtBox.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = text
        tp.font.size = Pt(15)
        tp.font.color.rgb = WHITE

    # ════════════════════════════════════════════
    # SLIDE 4 - Architecture
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Architecture Technique")

    pipeline_steps = [
        ("Environment", CYAN),
        ("Normal\nActivity", PURPLE),
        ("Attack\nEngine", RED),
        ("Telemetry", YELLOW),
        ("Detection", GREEN),
        ("Scoring", CYAN),
        ("AI Analyst", PURPLE),
        ("Report", GREEN),
    ]

    step_w = Inches(1.2)
    step_h = Inches(0.9)
    arrow_w = Inches(0.25)
    total_w = len(pipeline_steps) * step_w + (len(pipeline_steps) - 1) * arrow_w
    start_x = (SLIDE_WIDTH - total_w) // 2
    y = Inches(2.5)

    for i, (label, color) in enumerate(pipeline_steps):
        x = start_x + i * (step_w + arrow_w)
        rect = add_shape_rect(slide, x, y, step_w, step_h, BG_SLIGHTLY_LIGHTER, color)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = True
        rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < len(pipeline_steps) - 1:
            ax = x + step_w
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(0.3),
                                            arrow_w, Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_GRAY
            arrow.line.fill.background()

    # Bottom tech stack
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Backend: Python / FastAPI   |   Frontend: React   |   API: 28 endpoints"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 5 - Technologies
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Stack Technologique")

    # Backend column
    col1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5),
                           BG_SLIGHTLY_LIGHTER, CYAN)
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(5.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Backend"
    p.font.size = Pt(24)
    p.font.color.rgb = CYAN
    p.font.bold = True

    backend_items = ["Python 3.12", "FastAPI", "SQLite", "PyJWT", "Pydantic"]
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(4.5), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(backend_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Frontend column
    col2 = add_shape_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.5),
                           BG_SLIGHTLY_LIGHTER, PURPLE)
    txBox3 = slide.shapes.add_textbox(Inches(7.4), Inches(1.9), Inches(5.0), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Frontend"
    p3.font.size = Pt(24)
    p3.font.color.rgb = PURPLE
    p3.font.bold = True

    frontend_items = ["React 18", "Vite", "Tailwind CSS", "Recharts"]
    txBox4 = slide.shapes.add_textbox(Inches(7.7), Inches(2.6), Inches(4.5), Inches(2.5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    for i, item in enumerate(frontend_items):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

    # Bottom bar
    txBox5 = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "MITRE ATT&CK   |   Windows Event IDs   |   Sysmon"
    p5.font.size = Pt(18)
    p5.font.color.rgb = YELLOW
    p5.font.bold = True
    p5.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 6 - Jumeau Numerique
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Le Jumeau Numerique")

    items = [
        ("\u25b6  7 hosts simules dans l'environnement virtuel", WHITE),
        ("\u25b6  3 segments reseau : User LAN, Server Zone, DMZ", WHITE),
        ("\u25b6  5 utilisateurs avec profils realistes", WHITE),
        ("", WHITE),
        ("Activite normale simulee :", CYAN),
        ("    \u2022  Login / Logout", LIGHT_GRAY),
        ("    \u2022  Email (envoi, reception)", LIGHT_GRAY),
        ("    \u2022  Navigation Web & DNS", LIGHT_GRAY),
        ("    \u2022  Operations fichiers", LIGHT_GRAY),
        ("    \u2022  Requetes base de donnees", LIGHT_GRAY),
    ]
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20)
        p.font.color.rgb = color
        if color == CYAN:
            p.font.bold = True
        p.space_before = Pt(8)

    # ════════════════════════════════════════════
    # SLIDE 7 - Scenarios d'Attaque
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "4 Scenarios Bases sur des Menaces Reelles")

    rows = [
        ["Scenario", "Threat Actor", "Severite", "Techniques"],
        ["Spear Phishing", "APT29 Cozy Bear", "Critical", "6 phases"],
        ["Brute Force SSH", "TeamTNT", "High", "6 phases"],
        ["Lateral Movement", "APT28 Fancy Bear", "Critical", "6 phases"],
        ["Data Exfiltration", "Insider Threat", "High", "6 phases"],
    ]
    col_widths = [Inches(3.2), Inches(3.2), Inches(2.0), Inches(2.0)]
    add_table(slide, rows, col_widths, Inches(1.3), Inches(2.2), Inches(0.6))

    # ════════════════════════════════════════════
    # SLIDE 8 - Moteur de Detection
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    bullets = [
        "\u25b6  34 regles de detection implementees",
        "\u25b6  Fenetres glissantes temporelles",
        "\u25b6  Correlation d'incidents automatique",
        "",
        "Windows Event IDs :",
        "    4624, 4625, 4688, 4720, 5156 ...",
        "",
        "Sysmon Event IDs :",
        "    1, 3, 11, 22, 23 ...",
    ]
    add_bullet_slide(slide, "Moteur de Detection", bullets)

    # ════════════════════════════════════════════
    # SLIDE 9 - Scoring
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Scoring Multi-Dimensionnel")

    dimensions = [
        ("Detection (35%)", "Phases d'attaque detectees", CYAN),
        ("Couverture (30%)", "Techniques MITRE couvertes", PURPLE),
        ("Visibilite (20%)", "Sources de logs actives", GREEN),
        ("Reponse (15%)", "Temps moyen de detection", YELLOW),
    ]

    for i, (dim, desc, color) in enumerate(dimensions):
        y = Inches(1.9) + i * Inches(0.95)
        rect = add_shape_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.8),
                               BG_SLIGHTLY_LIGHTER, color)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{dim}  -  {desc}"
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

    # Maturity levels
    levels = ["Initial", "Repetable", "Defini", "Gere", "Optimise"]
    level_colors = [RED, YELLOW, CYAN, PURPLE, GREEN]
    lw = Inches(1.8)
    lh = Inches(0.7)
    start_x = Inches(1.5)
    y = Inches(5.8)

    txLabel = slide.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(6.0), Inches(0.4))
    lp = txLabel.text_frame.paragraphs[0]
    lp.text = "5 niveaux de maturite :"
    lp.font.size = Pt(16)
    lp.font.color.rgb = LIGHT_GRAY

    for i, (level, lc) in enumerate(zip(levels, level_colors)):
        x = start_x + i * (lw + Inches(0.15))
        rect = add_shape_rect(slide, x, y, lw, lh, BG_SLIGHTLY_LIGHTER, lc)
        tf = rect.text_frame
        p = tf.paragraphs[0]
        p.text = level
        p.font.size = Pt(14)
        p.font.color.rgb = lc
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 10 - Resultats
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Resultats de Detection")

    rows = [
        ["Scenario", "Score", "Detection", "Couverture", "Visibilite"],
        ["APT29 Phishing", "89.2 / 100", "83%", "83%", "100%"],
        ["TeamTNT Brute Force", "89.2 / 100", "83%", "83%", "100%"],
        ["APT28 Lateral Mvt", "87.5 / 100", "83%", "83%", "92%"],
        ["Insider Exfil", "78.3 / 100", "67%", "67%", "100%"],
    ]
    col_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    add_table(slide, rows, col_widths, Inches(0.8), Inches(2.2), Inches(0.65))

    # ════════════════════════════════════════════
    # SLIDE 11 - AI Analyst
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    bullets = [
        "\u25b6  NLG Rule-Based (pas d'API externe requise)",
        "\u25b6  Narrative executive automatique",
        "\u25b6  Extraction d'IOCs (IPs, domains, hashes)",
        "\u25b6  Impact compliance (GDPR, ISO 27001, NIST CSF)",
        "\u25b6  Recommandations strategiques prioritisees",
    ]
    add_bullet_slide(slide, "Analyste AI Integre", bullets)

    # ════════════════════════════════════════════
    # SLIDE 12 - Interface
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Interface Utilisateur")

    ui_pages_left = [
        "Dashboard avec KPIs animes",
        "MITRE ATT&CK Heatmap",
        "Threat Intelligence Feed",
        "Timeline interactive",
        "Export PDF professionnel",
        "Simulation Live (WebSocket)",
    ]
    ui_pages_right = [
        "Scenarios Manager",
        "Detection Rules Viewer",
        "Incidents & Alerts",
        "AI Analysis Reports",
        "Environment Overview",
        "Settings & Configuration",
    ]

    col1 = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.5), Inches(5.0))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    for i, item in enumerate(ui_pages_left):
        if i == 0:
            p = tf1.paragraphs[0]
        else:
            p = tf1.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(19)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)

    col2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(ui_pages_right):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(19)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)

    # ════════════════════════════════════════════
    # SLIDE 13 - Tests
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    bullets = [
        "\u2705  100 tests unitaires - 100% pass",
        "\u2705  8 fichiers de test couvrant tous les modules",
        "\u2705  pytest + pytest-cov",
        "\u2705  Docker Compose pour le deploiement",
    ]
    add_bullet_slide(slide, "Qualite & Tests", bullets)

    # ════════════════════════════════════════════
    # SLIDE 14 - Chiffres Cles
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Le Projet en Chiffres")

    stats = [
        ("11,315", "lignes de code", CYAN),
        ("13", "modules backend", PURPLE),
        ("12", "pages frontend", GREEN),
        ("34", "regles de detection", RED),
        ("100", "tests unitaires", YELLOW),
        ("28", "endpoints API", CYAN),
        ("4", "scenarios reels", PURPLE),
    ]

    # 2 rows layout: 4 top, 3 bottom
    box_w = Inches(2.5)
    box_h = Inches(1.8)
    gap = Inches(0.35)

    for i, (number, label, color) in enumerate(stats):
        if i < 4:
            row = 0
            col = i
            row_count = 4
        else:
            row = 1
            col = i - 4
            row_count = 3

        row_total_w = row_count * box_w + (row_count - 1) * gap
        row_start_x = (SLIDE_WIDTH - row_total_w) // 2
        x = row_start_x + col * (box_w + gap)
        y = Inches(1.8) + row * (box_h + Inches(0.3))

        rect = add_shape_rect(slide, x, y, box_w, box_h, BG_SLIGHTLY_LIGHTER, color)

        # Number
        numBox = slide.shapes.add_textbox(x, y + Inches(0.2), box_w, Inches(0.8))
        ntf = numBox.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = number
        np_.font.size = Pt(38)
        np_.font.color.rgb = color
        np_.font.bold = True
        np_.alignment = PP_ALIGN.CENTER

        # Label
        lblBox = slide.shapes.add_textbox(x, y + Inches(1.0), box_w, Inches(0.6))
        ltf = lblBox.text_frame
        lp_ = ltf.paragraphs[0]
        lp_.text = label
        lp_.font.size = Pt(15)
        lp_.font.color.rgb = LIGHT_GRAY
        lp_.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 15 - Demo
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Demonstration"
    p.font.size = Pt(48)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Demo live de la plateforme CyberTwin SOC"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.3), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "http://localhost:3001"
    p3.font.size = Pt(20)
    p3.font.color.rgb = PURPLE
    p3.alignment = PP_ALIGN.CENTER

    # ════════════════════════════════════════════
    # SLIDE 16 - Conclusion
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_title(slide, "Conclusion & Perspectives")

    # Apports
    apports_box = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                                  BG_SLIGHTLY_LIGHTER, GREEN)
    txH = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(5.0), Inches(0.5))
    hp = txH.text_frame.paragraphs[0]
    hp.text = "Apports"
    hp.font.size = Pt(24)
    hp.font.color.rgb = GREEN
    hp.font.bold = True

    apports = [
        "Approche Digital Twin innovante\npour la cybersecurite",
        "Simulation realiste basee sur des\nmenaces documentees",
        "Scoring objectif et\nmulti-dimensionnel",
    ]
    txA = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(4.5), Inches(3.5))
    tfa = txA.text_frame
    tfa.word_wrap = True
    for i, item in enumerate(apports):
        if i == 0:
            p = tfa.paragraphs[0]
        else:
            p = tfa.add_paragraph()
        p.text = f"\u2713  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)

    # Perspectives
    persp_box = add_shape_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5),
                                BG_SLIGHTLY_LIGHTER, CYAN)
    txH2 = slide.shapes.add_textbox(Inches(7.4), Inches(1.9), Inches(5.0), Inches(0.5))
    hp2 = txH2.text_frame.paragraphs[0]
    hp2.text = "Perspectives"
    hp2.font.size = Pt(24)
    hp2.font.color.rgb = CYAN
    hp2.font.bold = True

    perspectives = [
        "Integration de feeds Threat\nIntelligence en temps reel",
        "Support de scenarios\npersonnalises avances",
        "Deploiement cloud\n(AWS / Azure)",
        "Integration avec de vrais SIEM\n(Splunk, Sentinel)",
    ]
    txP = slide.shapes.add_textbox(Inches(7.7), Inches(2.6), Inches(4.5), Inches(3.5))
    tfp = txP.text_frame
    tfp.word_wrap = True
    for i, item in enumerate(perspectives):
        if i == 0:
            p = tfp.paragraphs[0]
        else:
            p = tfp.add_paragraph()
        p.text = f"\u2192  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)

    # ════════════════════════════════════════════
    # SLIDE 17 - Merci
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Merci"
    p.font.size = Pt(54)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Questions ?"
    p2.font.size = Pt(32)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # ── Save ──
    output_path = "CyberTwin_SOC_Soutenance.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
